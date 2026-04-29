"""
Create a presentation on refusal directions research extension to Hindi and Bengali.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 78, 121)  # Dark blue
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, content_list):
    """Add content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 121)
    
    # Add line under title
    line_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    line_shape.line.color.rgb = RGBColor(31, 78, 121)
    line_shape.line.width = Pt(2)
    
    # Content
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(5.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        
        if isinstance(item, tuple):
            main_text, sub_items = item
            p.text = main_text
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(31, 78, 121)
            p.level = 0
            
            for sub_item in sub_items:
                text_frame.add_paragraph()
                sub_p = text_frame.paragraphs[-1]
                sub_p.text = sub_item
                sub_p.font.size = Pt(16)
                sub_p.font.color.rgb = RGBColor(50, 50, 50)
                sub_p.level = 1
        else:
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.level = 0
    
    return slide


# Slide 1: Title
add_title_slide(
    prs,
    "Extending Refusal Directions to Indian Languages",
    "Improving Safety Evaluation for Hindi & Bengali LLMs\nDS 606 Project"
)

# Slide 2: Background - What is the paper about? (Part 1)
add_content_slide(
    prs,
    "Paper Overview: Refusal Directions",
    [
        ("Problem:", [
            "LLMs can be jailbroken by manipulating internal representations",
            "Safety mechanisms have identifiable 'refusal directions' in activation space"
        ]),
        ("Method:", [
            "Extract refusal vector from difference-in-means between safe/harmful completions",
            "Ablate (remove) this vector to measure jailbreak effectiveness",
            "Layer 16 identified as critical for multilingual safety"
        ])
    ]
)

# Slide 3: Background - What is the paper about? (Part 2)
add_content_slide(
    prs,
    "Key Findings from Original Research",
    [
        ("Attack Success Rate (ASR) Measurement:", [
            "Current method: Substring matching (looks for 'I cannot', 'I'm sorry', etc.)",
            "Simple but fragile—rephrased refusals easily missed",
            "⚠️ Not robust for cross-lingual evaluation"
        ]),
        ("Multilingual Gap:", [
            "Most research focuses on English LLMs",
            "Indian languages (Hindi, Bengali) largely unexplored",
            "Different linguistic structures may have different refusal mechanisms"
        ])
    ]
)

# Slide 4: What's Missing - Indian Languages Context
add_content_slide(
    prs,
    "Gap: Indian Languages Evaluation",
    [
        "Limited research on safety mechanisms in Hindi & Bengali LLMs",
        ("Why it matters:", [
            "2B+ speakers of Hindi and Bengali globally",
            "Different grammar, sentence structure, semantic patterns",
            "Refusal tokens and directions may differ from English"
        ]),
        ("Our Dataset:", [
            "Translated harmful ~100 prompts to Hindi & Bengali from PolyRefuse",
            "Identified language-specific refusal tokens through sampling",
            "Initial ASR: Hindi 98.4% | Bengali 95.3% (with jailbreak)"
        ])
    ]
)

# Slide 5: Our Contributions
add_content_slide(
    prs,
    "Our Contribution: What We're Doing",
    [
        ("Extension to Indian Languages:", [
            "✓ Hindi & Bengali dataset integration",
            "✓ Refusal token identification & filtering",
            "✓ Ablation pipeline execution (both languages)"
        ]),
        ("Robust Evaluation Framework:", [
            "Current: Substring matching (fragile, language-dependent)",
            "Proposed: Llama Guard 4 (MLCommons taxonomy)",
            "✓ Classifies responses to 14 hazard categories",
            "✓ Multilingual support built-in"
        ])
    ]
)

# Slide 6: ASR Evaluation - Current vs Proposed
add_content_slide(
    prs,
    "Attack Success Rate: Old vs New Evaluation",
    [
        ("Current Method (Substring Matching):", [
            "❌ Checks only for exact keywords: 'sorry', 'cannot', 'unable'",
            "❌ Misses semantic refusals with different phrasing",
            "❌ No semantic understanding",
            "⚠️ Results: Hindi 98.4%, Bengali 95.3%"
        ]),
        ("Proposed Method (Llama Guard 4):", [
            "✓ Uses Llama Guard 4-12B safety classifier",
            "✓ Understands semantic safety violations (14 MLCommons categories)",
            "✓ Robust to rephrasing and language variations",
            "→ More reliable cross-lingual evaluation"
        ])
    ]
)

# Slide 7: Next Steps & Timeline
add_content_slide(
    prs,
    "Next Steps & Timeline",
    [
        ("Immediate (This Week):", [
            "✓ Deploy Llama Guard 4 evaluation on server",
            "Generate robust ASR metrics for Hindi & Bengali",
            "Compare substring vs Llama Guard 4 results"
        ]),
        ("Short-term (Next 2 weeks):", [
            "Statistical analysis: robustness improvements",
            "Visualization & comparison plots",
            "Write-up for paper/presentation"
        ]),
        ("Deliverables:", [
            "Paper with Indian language experiments",
            "Public dataset: Hindi/Bengali harmful prompts & refusals",
            "Code: Llama Guard 4 evaluation pipeline"
        ])
    ]
)

# Save
output_file = "/Users/sravani/Documents/VSCode_projects/DS606_project_r2/Refusal_Directions_Indian_Languages_Presentation.pptx"
prs.save(output_file)
print(f"✓ Presentation created: {output_file}")
